

from pathlib import Path
from typing import List, Dict
import logging

from sentence_transformers import SentenceTransformer

# Document parsing libraries
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from .utils_textclean import clean_pdf_text, suggest_ocr_needed, detect_and_extract_schedule_from_pdf
except ImportError:
    try:
        from utils_textclean import clean_pdf_text, suggest_ocr_needed, detect_and_extract_schedule_from_pdf
    except ImportError:
        # Fallback if text cleaning module is not available
        def clean_pdf_text(text):
            return text

        def suggest_ocr_needed(path, text):
            return False

        def detect_and_extract_schedule_from_pdf(path, text):
            return ""

# Multilingual E5 model for better Spanish support
MODEL_NAME = "intfloat/multilingual-e5-base"
E5_PREFIX_Q = "query: "
E5_PREFIX_P = "passage: "
_model = None


def get_model():
    global _model
    if _model is None:
        _model = SentenceTransformer(MODEL_NAME)
    return _model


def expand_query(query: str) -> str:
    """
    Expand Spanish queries with synonyms and common variations
    to improve retrieval on short queries like "¿Cuáles son las nubes?"
    """
    expansion_dict = {
        "nube": ["cloud", "clouds", "plataforma en la nube", "Azure", "AWS", "MongoDB Atlas"],
        "nubes": ["cloud", "clouds", "plataformas en la nube", "Azure", "AWS", "MongoDB Atlas"],
        "entrega": ["fecha de entrega", "fecha límite", "deadline"],
        "evaluación": ["puntaje", "porcentaje", "evaluaciones", "nota"],
        "evaluaciones": ["puntaje", "porcentaje", "evaluación", "nota"],
        "examen": ["evaluación", "prueba", "test"],
        "horario": ["cronograma", "calendario", "fechas"],
        "profesor": ["docente", "instructor", "maestro"],
    }

    ql = query.lower()
    extra_terms = []

    for key, synonyms in expansion_dict.items():
        if key in ql:
            extra_terms.extend(synonyms)

    if extra_terms:
        # Add unique terms to query
        unique_extras = list(set(extra_terms))
        # Limit to 5 extra terms
        return query + " " + " ".join(unique_extras[:5])

    return query


def embed_e5(texts, is_query: bool):
    """
    E5 embeddings with proper prefixes for better Spanish support:
      - queries:  'query: ...'
      - passages: 'passage: ...'

    Args:
        texts: List of strings to embed
        is_query: True for query mode (adds query prefix), False for passage mode

    Returns:
        List of normalized embeddings (768 dimensions for e5-base)
    """
    model = get_model()
    prefix = E5_PREFIX_Q if is_query else E5_PREFIX_P
    prefixed = [prefix + t.strip() for t in texts]
    # Keep normalized for cosine similarity
    vecs = model.encode(prefixed, normalize_embeddings=True)
    return vecs.tolist()


def extract_text_from_pdf(file_path: Path) -> tuple[str, str]:
    """
    Extract text from PDF file using multiple libraries with AES decryption support

    Returns:
        tuple: (text, schedule_table)
            - text: cleaned text from the PDF
            - schedule_table: extracted schedule table if found, empty string otherwise
    """

    # Common passwords to try for encrypted PDFs
    file_path_obj = Path(file_path)
    common_passwords = ["", "password", "123456",
                        "admin", "user", file_path_obj.stem.lower()]

    # Try pdfplumber first (best for complex layouts)
    if pdfplumber is not None:
        try:
            with pdfplumber.open(file_path, password=None) as pdf:
                text = ""
                raw_full_text = ""  # Collect raw text for table extraction
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        raw_full_text += page_text + "\n\n"  # Keep raw for table extraction
                        # Clean the extracted text
                        cleaned_text = clean_pdf_text(page_text)
                        if cleaned_text.strip():
                            text += f"Page {page_num + 1}:\n{cleaned_text}\n\n"

                if text.strip():
                    # Try to extract schedule table from raw text
                    schedule_content = detect_and_extract_schedule_from_pdf(
                        str(file_path), raw_full_text)
                    return text.strip(), schedule_content
                return "", ""
        except Exception as e:
            if "password" not in str(e).lower():
                print(f"⚠️  pdfplumber failed for {file_path_obj.name}: {e}")

    # Try PyMuPDF second (good for most PDFs)
    if fitz is not None:
        try:
            doc = fitz.open(str(file_path))
            text = ""
            raw_full_text = ""  # Collect raw text for table extraction
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                if page_text.strip():
                    raw_full_text += page_text + "\n\n"  # Keep raw for table extraction
                    # Clean the extracted text
                    cleaned_text = clean_pdf_text(page_text)
                    if cleaned_text.strip():
                        text += f"Page {page_num + 1}:\n{cleaned_text}\n\n"
            doc.close()
            if text.strip():
                # Try to extract schedule table from raw text
                schedule_content = detect_and_extract_schedule_from_pdf(
                    str(file_path), raw_full_text)
                return text.strip(), schedule_content
            return "", ""
        except Exception as e:
            if "password" not in str(e).lower():
                print(f"⚠️  PyMuPDF failed for {file_path_obj.name}: {e}")

    # Enhanced PyPDF2 with password attempts and AES support
    if PyPDF2 is None:
        logging.warning(f"No PDF libraries available, skipping {file_path}")
        return "", ""

    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""

            # Handle password-protected PDFs with multiple password attempts
            if reader.is_encrypted:
                decrypted = False
                for password in common_passwords:
                    try:
                        if reader.decrypt(password):
                            print(
                                f"🔓 Successfully decrypted {file_path_obj.name} with password: {'(empty)' if password == '' else '***'}")
                            decrypted = True
                            break
                    except Exception as decrypt_error:
                        continue

                if not decrypted:
                    print(
                        f"🔒 Could not decrypt {file_path_obj.name} - password protected")
                    print(
                        f"   Tried passwords: {len(common_passwords)} common passwords")
                    return "", ""

            # Extract text from all pages
            raw_full_text = ""  # Collect raw text for table extraction
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        raw_full_text += page_text + "\n\n"  # Keep raw for table extraction
                        # Clean the extracted text
                        cleaned_text = clean_pdf_text(page_text)
                        if cleaned_text.strip():
                            text += f"Page {page_num + 1}:\n{cleaned_text}\n\n"
                except Exception as e:
                    print(
                        f"⚠️  Could not extract text from page {page_num + 1} in {file_path_obj.name}: {e}")
                    continue

            if text.strip():
                # Check if text quality is too poor and suggest OCR
                if suggest_ocr_needed(str(file_path), text):
                    print(f"💡 Consider running OCR on {file_path_obj.name}:")
                    print(
                        f"   ocrmypdf --force-ocr --skip-text --optimize 1 '{file_path}' '{file_path.stem}.ocr.pdf'")

                # Automatically detect and extract schedule table from course guide PDFs
                schedule_content = detect_and_extract_schedule_from_pdf(
                    str(file_path), raw_full_text)
                return text.strip(), schedule_content
            else:
                print(
                    f"⚠️  No readable text found in {file_path_obj.name} (might be image-based PDF)")
                print(
                    f"💡 Try OCR: ocrmypdf --force-ocr '{file_path}' '{file_path.stem}.ocr.pdf'")

                return "", ""

    except Exception as e:
        error_msg = str(e)
        if "PyCryptodome" in error_msg:
            print(
                f"🔐 {file_path_obj.name} requires advanced encryption - trying alternative method...")
            # The libraries are now installed, so this should work
            return "", ""
        else:
            logging.error(f"Error reading PDF {file_path}: {e}")
            print(f"❌ PDF extraction failed for {file_path_obj.name}: {e}")
            return "", ""


def extract_text_from_docx(file_path: Path) -> str:
    """Extract text from Word document"""
    if Document is None:
        logging.warning(f"python-docx not installed, skipping {file_path}")
        return ""

    try:
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return "\n".join(text)
    except Exception as e:
        logging.error(f"Error reading DOCX {file_path}: {e}")
        return ""


def extract_text_from_excel(file_path: Path) -> str:
    """Extract text from Excel file"""
    if load_workbook is None:
        logging.warning(f"openpyxl not installed, skipping {file_path}")
        return ""

    try:
        workbook = load_workbook(file_path, data_only=True)
        text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text.append(f"Sheet: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                row_text = [
                    str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_text):  # Skip empty rows
                    text.append(" | ".join(row_text))
        return "\n".join(text)
    except Exception as e:
        logging.error(f"Error reading Excel {file_path}: {e}")
        return ""


def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from PowerPoint presentation"""
    if Presentation is None:
        logging.warning(f"python-pptx not installed, skipping {file_path}")
        return ""

    try:
        prs = Presentation(file_path)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"Slide {slide_num}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logging.error(f"Error reading PPTX {file_path}: {e}")
        return ""


def extract_text_from_html(file_path: Path) -> str:
    """Extract text from HTML file"""
    if BeautifulSoup is None:
        logging.warning(f"beautifulsoup4 not installed, skipping {file_path}")
        return ""

    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            soup = BeautifulSoup(file.read(), 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        logging.error(f"Error reading HTML {file_path}: {e}")
        return ""


def read_texts(folder: str) -> List[Dict]:
    """Read text from various document formats"""
    items = []
    folder_path = Path(folder)


def extract_text_from_csv(file_path: Path) -> str:
    """Extract text from CSV file"""
    try:
        import csv
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            # Read CSV and convert to text format
            reader = csv.reader(file)
            text = []
            for row_num, row in enumerate(reader):
                if row_num == 0:  # Header row
                    text.append("Headers: " + " | ".join(row))
                else:
                    text.append(f"Row {row_num}: " + " | ".join(row))
            return "\n".join(text)
    except Exception as e:
        logging.error(f"Error reading CSV {file_path}: {e}")
        return ""


def read_texts(folder: str) -> List[Dict]:
    """Read text from various document formats"""
    items = []
    folder_path = Path(folder)

    # Define file type mappings
    extractors = {
        '.txt': lambda p: p.read_text(encoding="utf-8", errors="ignore"),
        '.md': lambda p: p.read_text(encoding="utf-8", errors="ignore"),
        '.csv': extract_text_from_csv,
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.xlsx': extract_text_from_excel,
        '.xls': extract_text_from_excel,
        '.pptx': extract_text_from_pptx,
        '.html': extract_text_from_html,
        '.htm': extract_text_from_html,
    }

    # Process all supported file types
    for extension, extractor in extractors.items():
        for file_path in folder_path.glob(f"**/*{extension}"):
            try:
                # Handle PDF files which return tuple (text, schedule)
                if extension == '.pdf':
                    text, schedule_table = extractor(file_path)
                    # Store schedule in metadata if present
                    if schedule_table:
                        items.append({
                            "path": str(file_path),
                            "text": text,
                            "file_type": extension,
                            "schedule": schedule_table
                        })
                    else:
                        items.append({
                            "path": str(file_path),
                            "text": text,
                            "file_type": extension
                        })
                else:
                    text = extractor(file_path)
                    if text.strip():  # Only add if text was extracted
                        items.append({
                            "path": str(file_path),
                            "text": text,
                            "file_type": extension
                        })

                if text.strip():
                    print(f"✅ Processed: {file_path.name} ({extension})")
                else:
                    print(f"⚠️  No text extracted from: {file_path.name}")
            except Exception as e:
                print(f"❌ Error processing {file_path.name}: {e}")
                logging.error(f"Error processing {file_path}: {e}")

    print(f"\n📊 Total documents processed: {len(items)}")
    return items


def chunk_text(text: str, max_tokens: int = 220, overlap_tokens: int = 50) -> List[str]:
    """
    Smart chunking with overlap for better context preservation.
    Uses word-based tokens for simplicity.

    Args:
        text: Text to chunk
        max_tokens: Maximum words per chunk (220 ≈ 150-220 words)
        overlap_tokens: Overlap between chunks (50 ≈ 20-30% overlap)

    Returns:
        List of text chunks
    """
    import re

    # Additional cleaning before chunking
    text = clean_chunk_text(text)

    # Simple whitespace token proxy
    words = text.split()

    if len(words) <= max_tokens:
        return [text] if text.strip() else []

    chunks = []
    i = 0

    while i < len(words):
        chunk_words = words[i:i + max_tokens]
        chunk = " ".join(chunk_words)

        # Drop boilerplate like "Page 1:" if it's at the start
        if chunk.startswith("Page ") and ":" in chunk[:12]:
            chunk = chunk.split(":", 1)[1].strip()

        if chunk.strip():
            chunks.append(chunk)

        # Move forward with overlap
        i += max_tokens - overlap_tokens

    return chunks


def clean_chunk_text(text: str) -> str:
    """
    Clean text before chunking to improve retrieval quality.
    Removes PDF artifacts, normalizes spacing, and fixes common OCR errors.
    """
    import re

    if not text or not text.strip():
        return ""

    # Remove soft hyphens
    text = text.replace('\u00ad', '')

    # Merge hyphenated words at line breaks: "evalua-\nción" → "evaluación"
    text = re.sub(r'-\s*\n\s*', '', text)

    # Trim spaces before newlines
    text = re.sub(r'\s+\n', '\n', text)

    # Collapse multiple blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)

    # Collapse multiple spaces/tabs
    text = re.sub(r'[ \t]{2,}', ' ', text)

    # Remove page numbers that appear as standalone patterns
    text = re.sub(r'\n\d+\n', '\n', text)

    return text.strip()


def make_chunk_header(path: str, page_num: int = None) -> str:
    """
    Create a stable header for each chunk to improve retrieval.
    Example: "[DOC:Guia de Curso.pdf] [PAGE:3] "

    This helps queries like "which clouds will we use?" match chunks
    containing "Nube usando Azure SQL… MongoDB Atlas… (Amazon Data Lake)…"
    """
    from pathlib import Path

    doc_name = Path(path).name
    header = f"[DOC:{doc_name}]"

    if page_num is not None:
        header += f" [PAGE:{page_num}]"

    return header + " "


def embed_texts(texts: List[str]):
    """Legacy function - use embed_e5 for better results"""
    model = get_model()
    return model.encode(texts, normalize_embeddings=True).tolist()
